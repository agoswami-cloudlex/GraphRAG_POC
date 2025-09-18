import os
from typing import Optional

from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

def extract_text_from_pdf(file_path: str) -> str:
    text = ""
    with open(file_path, "rb") as f:
        reader = PdfReader(f)
        for page in reader.pages:
            text += page.extract_text() or ""
    return text

def extract_text_from_docx(file_path: str) -> str:
    doc = DocxDocument(file_path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def extract_text_from_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()

def extract_text(file_path: str) -> Optional[str]:
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".pdf":
            return extract_text_from_pdf(file_path)
        elif ext == ".docx":
            return extract_text_from_docx(file_path)
        elif ext == ".pptx":
            return extract_text_from_pptx(file_path)
        elif ext in [".txt", ".md"]:
            return extract_text_from_txt(file_path)
        else:
            return None
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}")
        return None
